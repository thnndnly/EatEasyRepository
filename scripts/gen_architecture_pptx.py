#!/usr/bin/env python3
"""Baut EatEasy-Architektur.pptx aus den gerenderten PDF-Folien.

Die Grafik-Folien werden in gen_architecture_pdf.py als Vektor-PDF gezeichnet
(Hausstil, geprueft). Hier rendern wir jede Seite hochaufgeloest und legen sie
als vollflaechiges Bild auf eine 16:9-Folie — so ist die PPTX pixelgleich zum
PDF, ohne auf die (in dieser Umgebung defekte) LibreOffice-Konvertierung zu
setzen.
"""
import os
import tempfile

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

# Repo-Wurzel relativ zum Skript (scripts/..) — plattformunabhaengig.
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PDF = os.path.join(ROOT, "EatEasy-Architektur.pdf")
OUT = os.path.join(ROOT, "EatEasy-Architektur.pptx")
TMP = tempfile.mkdtemp(prefix="eateasy_slides_")

doc = fitz.open(PDF)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for i, page in enumerate(doc):
    png = os.path.join(TMP, f"slide_{i + 1}.png")
    page.get_pixmap(dpi=200).save(png)  # ~2667x1500 px
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print("saved", OUT, "—", len(doc), "Folien")
